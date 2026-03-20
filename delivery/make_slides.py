"""
Generate EDAP Executive Slide Pack (4 slides)
Positions delivery scope against the broader enterprise wiki context.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # slide headers / strong text
TEAL    = RGBColor(0x00, 0x6E, 0x8C)   # accents, dividers
BLUE    = RGBColor(0x3A, 0x7B, 0xBF)   # secondary accent / box fills
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF7)   # light background panels
DGRAY   = RGBColor(0x2C, 0x2C, 0x2C)   # body text
MID     = RGBColor(0x55, 0x6B, 0x8A)   # muted label text

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_label(slide, text, left, top, width, height,
              font_size=Pt(11), bold=False, color=DGRAY,
              align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_para(tf, text, font_size=Pt(10), bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, space_before=Pt(3)):
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def header_band(slide, title, subtitle=None):
    """Dark navy header band across the top."""
    band = add_rect(slide, 0, 0, 10, 1.15, fill_rgb=NAVY)
    add_label(slide, title, 0.3, 0.12, 9.0, 0.55,
              font_size=Pt(22), bold=True, color=WHITE)
    if subtitle:
        add_label(slide, subtitle, 0.3, 0.68, 9.0, 0.38,
                  font_size=Pt(11), color=RGBColor(0xB0, 0xC8, 0xE0), italic=False)
    # teal rule
    add_rect(slide, 0, 1.15, 10, 0.06, fill_rgb=TEAL)


def footer(slide, page_num, total=4):
    add_rect(slide, 0, 7.25, 10, 0.25, fill_rgb=NAVY)
    add_label(slide, "Water Corporation  |  EDAP Programme  |  March 2026",
              0.3, 7.26, 7.0, 0.22, font_size=Pt(8), color=RGBColor(0xB0, 0xC8, 0xE0))
    add_label(slide, f"{page_num} / {total}", 9.3, 7.26, 0.6, 0.22,
              font_size=Pt(8), color=RGBColor(0xB0, 0xC8, 0xE0), align=PP_ALIGN.RIGHT)


def bullet_box(slide, title, bullets, left, top, width, height,
               title_color=TEAL, body_size=Pt(9.5)):
    """A titled bullet box with light grey background."""
    add_rect(slide, left, top, width, height, fill_rgb=LGRAY,
             line_rgb=RGBColor(0xCC, 0xD6, 0xE0), line_width=Pt(0.75))
    add_label(slide, title, left + 0.12, top + 0.08, width - 0.2, 0.28,
              font_size=Pt(10), bold=True, color=title_color)
    # divider
    add_rect(slide, left + 0.12, top + 0.36, width - 0.24, 0.025, fill_rgb=TEAL)
    tb = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.42),
        Inches(width - 0.24), Inches(height - 0.55)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = body_size
        run.font.color.rgb = DGRAY


# ── Presentation setup ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Enterprise Imperative
# ════════════════════════════════════════════════════════════════════════════

s1 = prs.slides.add_slide(blank_layout)
header_band(s1,
    "The Enterprise Imperative",
    "Why EDAP exists — and what it changes for Water Corporation")
footer(s1, 1)

# Intro sentence
add_label(s1,
    "Water Corporation's data landscape has grown organically across seven business domains, "
    "creating fragmentation that undermines trust, slows decisions, and complicates compliance.",
    0.3, 1.32, 9.4, 0.45, font_size=Pt(10.5), color=DGRAY)

# Two columns: Current State | Target State
# Left panel
bullet_box(s1, "Current State — Fragmented by Design", [
    "Six source systems (SAP, Maximo, SCADA, GIS, Salesforce, CASS) in silos",
    "Multiple incompatible pipeline technologies: Glue, ADF, SAP DS, SageMaker",
    "No unified governance — classification inconsistent, access ad hoc",
    "Knowledge concentrated in individuals, not embedded in platform",
    "Regulatory obligations (PRIS Act, SOCI Act, WAICP) met manually, not systematically",
], 0.3, 1.87, 4.55, 2.8)

# Right panel
bullet_box(s1, "Target State — Governed by Architecture", [
    "Single platform: Databricks + Unity Catalog across all seven domains",
    "Metadata-driven pipelines — new source onboarded in days, not weeks",
    "Self-governing data: WAICP classification enforced at column level via ABAC",
    "Certified data products with contracts, SLAs, and quality scores",
    "AI/ML as a first-class citizen with model inventory and risk classification",
], 5.15, 1.87, 4.55, 2.8)

# Central arrow
add_rect(s1, 4.73, 2.85, 0.54, 0.45, fill_rgb=TEAL)
add_label(s1, "→", 4.73, 2.83, 0.54, 0.45,
          font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three principles bar at the bottom
bar_y = 4.82
bar_items = [
    ("Comply by Design",   "PRIS, SOCI, WAICP, State Records\nenforced architecturally — not manually"),
    ("Data as a Product",  "Findable, accessible, understandable,\nquality-assured, dependable"),
    ("Team Self-Sufficiency", "WC engineers own their domains;\nSI transfers capability, not dependency"),
]
bw = 3.0
for i, (title, body) in enumerate(bar_items):
    bx = 0.3 + i * (bw + 0.1)
    add_rect(s1, bx, bar_y, bw, 0.98, fill_rgb=NAVY,
             line_rgb=TEAL, line_width=Pt(1))
    add_label(s1, title, bx + 0.12, bar_y + 0.08, bw - 0.22, 0.28,
              font_size=Pt(9.5), bold=True, color=TEAL)
    add_label(s1, body, bx + 0.12, bar_y + 0.36, bw - 0.22, 0.56,
              font_size=Pt(8.5), color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Enterprise Architecture
# ════════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(blank_layout)
header_band(s2,
    "The Enterprise Architecture",
    "Medallion lakehouse with Unity Catalog governance — the foundation all delivery work builds on")
footer(s2, 2)

add_label(s2,
    "The wikis define the architecture that precedes delivery. Every scope item implements a layer of this stack.",
    0.3, 1.32, 9.4, 0.32, font_size=Pt(10.5), color=DGRAY)

# ── Medallion stack (left 5.5 inches) ──
stack_left = 0.3
stack_top  = 1.72
layer_w    = 5.3
layer_h    = 0.50
gap        = 0.06

layers = [
    (RGBColor(0xB8, 0x5C, 0x1A), "GOLD",   "Dimensional  |  BI Dashboards  |  Genie (conversational)  |  AI/BI"),
    (RGBColor(0xB8, 0x5C, 0x1A), "GOLD",   "Enriched  |  Exploratory  |  Sandbox"),
    (RGBColor(0x70, 0x8E, 0xB0), "SILVER", "Base  |  Protected (PI-masked)"),
    (RGBColor(0x70, 0x8E, 0xB0), "SILVER", "Processed  |  Raw"),
    (RGBColor(0x8B, 0x6E, 0x4A), "BRONZE", "Landing  |  Quarantine"),
]

for i, (col, tier, zones) in enumerate(layers):
    y = stack_top + i * (layer_h + gap)
    add_rect(s2, stack_left, y, layer_w, layer_h, fill_rgb=col)
    add_label(s2, tier, stack_left + 0.08, y + 0.06, 0.75, 0.35,
              font_size=Pt(8), bold=True, color=WHITE)
    add_label(s2, zones, stack_left + 0.9, y + 0.10, layer_w - 1.0, 0.30,
              font_size=Pt(9), color=WHITE)

add_label(s2, "Medallion Architecture  (10 zones)", stack_left, stack_top - 0.28,
          layer_w, 0.24, font_size=Pt(9), bold=True, color=TEAL)

# Governance overlay label
add_rect(s2, stack_left, stack_top + 5 * (layer_h + gap) + 0.04, layer_w, 0.30,
         fill_rgb=NAVY)
add_label(s2, "Unity Catalog  —  single metastore, domain-based catalogs, ABAC tag enforcement",
          stack_left + 0.1, stack_top + 5 * (layer_h + gap) + 0.07, layer_w - 0.15, 0.22,
          font_size=Pt(8.5), color=RGBColor(0xB0, 0xC8, 0xE0))

# ── Right panel: 3 governance pillars ──
rp_left = 5.9
rp_top  = 1.72

add_label(s2, "Governance Pillars", rp_left, rp_top - 0.28, 3.9, 0.24,
          font_size=Pt(9), bold=True, color=TEAL)

pillars = [
    ("4-Layer Tagging Model", [
        "Layer 1: WAICP Classification (mandatory)",
        "Layer 2: Sensitivity Reason (PII, SOCI, AI)",
        "Layer 3: Access & Handling (masking, contracts)",
        "Layer 4: Operational Metadata (owner, refresh, SLA)",
    ]),
    ("Access Control (ABAC)", [
        "Tag-driven row/column masking at query time",
        "Domain-federated ownership via MANAGE privilege",
        "Break-glass audit trail in Unity Catalog system tables",
    ]),
    ("Compliance by Architecture", [
        "PRIS Act 2024 — PI masking, anonymisation at source",
        "SOCI Act 2018 — CMK encryption, critical infra controls",
        "State Records Act — retention & disposal metadata",
    ]),
]

ph = 1.04
for i, (title, pts) in enumerate(pillars):
    py = rp_top + i * (ph + 0.06)
    bullet_box(s2, title, pts, rp_left, py, 3.82, ph, body_size=Pt(8.5))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Delivery Scope Positioned Against the Architecture
# ════════════════════════════════════════════════════════════════════════════

s3 = prs.slides.add_slide(blank_layout)
header_band(s3,
    "Nine Epics — Delivery Scope Mapped to the Architecture",
    "95 features, sequenced foundation-first, each implementing a layer of the enterprise design")
footer(s3, 3)

add_label(s3,
    "The nine epics are not independent workstreams — they are ordered architectural layers. "
    "This PI establishes the foundation that all subsequent capability depends on.",
    0.3, 1.32, 9.4, 0.38, font_size=Pt(10.5), color=DGRAY)

# Epic table — 3 columns: Epic | This PI focus | Wiki anchor
col_headers = ["Epic", "This PI Delivers", "Architectural Anchor"]
col_w = [2.8, 3.5, 3.2]
col_x = [0.3, 3.2, 6.8]
row_h = 0.44
tbl_top = 1.82

# Header row
for j, (hdr, w, x) in enumerate(zip(col_headers, col_w, col_x)):
    add_rect(s3, x, tbl_top, w - 0.05, 0.32, fill_rgb=NAVY)
    add_label(s3, hdr, x + 0.08, tbl_top + 0.04, w - 0.18, 0.24,
              font_size=Pt(9), bold=True, color=WHITE)

epics = [
    ("E1  Platform Foundation",
     "Environments, UC metastore, identity, compute policies, CI/CD (DABs)",
     "EDAP Access Model  |  Medallion Architecture"),
    ("E2  Data Ingestion & Integration",
     "Metadata-driven ingestion framework; CDC, batch & event patterns proven",
     "Pipeline Framework (EDAP-FWK-001)"),
    ("E4  Security, Governance & Compliance",
     "ABAC policies activated; 4-layer tagging applied; Alation OCF connector live",
     "Tagging Strategy  |  Data Governance Roles"),
    ("E5  DataOps & Engineering Excellence",
     "Branching strategy, testing pyramid, cost guardrails ≤10% non-prod",
     "Databricks End-to-End Platform"),
    ("E3/6/7  Transformation, Catalogue & AI",
     "Framework validated end-to-end; Alation–UC bidirectional sync proven",
     "Pipeline Framework  |  Medallion Architecture"),
    ("E8/9  Migration & Capability Uplift",
     "Migration pathways documented; WC engineers shadowing & self-sufficient",
     "Lifecycle Guides  |  EA Vision"),
]

row_colors = [LGRAY, WHITE]
for i, (epic, pi_focus, anchor) in enumerate(epics):
    ry = tbl_top + 0.32 + i * row_h + 0.02
    bg = row_colors[i % 2]
    for j, (w, x) in enumerate(zip(col_w, col_x)):
        add_rect(s3, x, ry, w - 0.05, row_h - 0.03,
                 fill_rgb=bg, line_rgb=RGBColor(0xCC, 0xD6, 0xE0), line_width=Pt(0.5))
    texts = [epic, pi_focus, anchor]
    bold_flags = [True, False, False]
    text_colors = [NAVY, DGRAY, MID]
    for j, (txt, w, x, bld, col) in enumerate(zip(texts, col_w, col_x, bold_flags, text_colors)):
        add_label(s3, txt, x + 0.08, ry + 0.04, w - 0.18, row_h - 0.1,
                  font_size=Pt(8.5), bold=bld, color=col)

# Bottom callout
add_rect(s3, 0.3, 7.0, 9.4, 0.20, fill_rgb=TEAL)
add_label(s3,
    "Scope items S2–S23 each map to one or more epics. The full feature set (95 features) is traceable to architectural decisions already closed in the wiki.",
    0.45, 7.01, 9.1, 0.18, font_size=Pt(8), color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How We Know It's Working
# ════════════════════════════════════════════════════════════════════════════

s4 = prs.slides.add_slide(blank_layout)
header_band(s4,
    "How We Know It's Working",
    "Success measures, risks, and the conditions for a self-sustaining platform")
footer(s4, 4)

add_label(s4,
    "The wikis define what 'done' looks like. The delivery scope is the path to get there. "
    "These measures confirm the architecture is landing — not just deployed.",
    0.3, 1.32, 9.4, 0.38, font_size=Pt(10.5), color=DGRAY)

# 4 measure boxes in 2×2 grid
measures = [
    ("Single Source of Truth",
     ["Any authorised user can find, understand, and access any data asset",
      "Alation + Unity Catalog show consistent lineage end-to-end",
      "No parallel pipelines serving the same domain"]),
    ("Compliance by Design",
     ["PRIS Act audit satisfiable from Unity Catalog system tables",
      "All PI columns masked or anonymised at Bronze → Silver boundary",
      "SOCI controls evidenced without manual intervention"]),
    ("Speed & Self-Service",
     ["New source system onboarded in days via configuration (not code)",
      "Domain teams run their own pipelines without platform team tickets",
      "Non-prod compute cost ≤10% of production"]),
    ("Trustworthy AI",
     ["All models have cards: training data, owner, risk tier, performance",
      "No model in production without an approved AI Risk Assessment",
      "WC engineers lead reviews and own model lifecycle"]),
]

gw = 4.55
gh = 2.1
for i, (title, pts) in enumerate(measures):
    gx = 0.3  + (i % 2) * (gw + 0.1)
    gy = 1.82 + (i // 2) * (gh + 0.12)
    bullet_box(s4, title, pts, gx, gy, gw, gh, body_size=Pt(9.5))

# Risk banner
add_rect(s4, 0.3, 6.15, 9.4, 0.25, fill_rgb=RGBColor(0x7A, 0x1F, 0x1F))
add_label(s4,
    "Key risk: scope without sequencing.  E1 Foundation must complete before E2–E7 can deliver at pace.  "
    "Any de-scoping of governance features (E4) breaks the compliance by design principle.",
    0.45, 6.16, 9.1, 0.22, font_size=Pt(8), color=WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────

out = "/home/user/DataWikis/delivery/edap-executive-slide-pack.pptx"
prs.save(out)
print(f"Saved: {out}")

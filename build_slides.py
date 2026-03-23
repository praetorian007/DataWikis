"""
Build EDAP Executive Slide Pack — 4 slides
  1. Portfolio Vision
  2. Program Strategic Themes
  3. Program Scope — Nine Epics
  4. Project Outcomes — Supporting the Vision
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide geometry ────────────────────────────────────────────────────────────
SW, SH = 9144000, 6858000   # 10" × 7.5"
HH     = 1051560             # header bar height
AH     = 54864               # teal accent line height
FT     = 6629400             # footer top
FH     = SH - FT             # footer height = 228600
CT     = HH + AH             # content top   = 1106424
MG     = 274320              # left/right margin
PM     = 109728              # inner padding

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
TEAL  = RGBColor(0x00, 0x6E, 0x8C)
LGRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x2C, 0x2C)
LBLUE = RGBColor(0xB0, 0xC8, 0xE0)
TEAL2 = RGBColor(0x00, 0x4E, 0x6A)   # darker teal variant

# ── Low-level helpers ─────────────────────────────────────────────────────────

def mk():
    p = Presentation()
    p.slide_width  = Emu(SW)
    p.slide_height = Emu(SH)
    return p

def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])

def bx(sl, l, t, w, h, c, border=False):
    s = sl.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = c
    if border:
        s.line.color.rgb = TEAL
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def tx(sl, l, t, w, h):
    tb = sl.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tb.text_frame.word_wrap = True
    return tb

def ln(tb_shape, text, pt, color,
       bold=False, italic=False,
       align=PP_ALIGN.LEFT,
       first=False,
       sb=0, sa=0):          # space_before / space_after in points
    tf = tb_shape.text_frame
    p  = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text         = text
    r.font.size    = Pt(pt)
    r.font.color.rgb = color
    r.font.bold    = bold
    r.font.italic  = italic
    return p

# ── Shared chrome (header + accent + footer) ─────────────────────────────────

def chrome(sl, title, subtitle, num):
    bx(sl, 0, 0, SW, HH, NAVY)
    bx(sl, 0, HH, SW, AH, TEAL)
    bx(sl, 0, FT, SW, SH - FT, NAVY)

    t = tx(sl, MG, PM, SW - 2*MG, 503000)
    ln(t, title, 22, WHITE, bold=True, first=True)

    s = tx(sl, MG, 619000, SW - 2*MG, 380000)
    ln(s, subtitle, 11, LBLUE, first=True)

    fl = tx(sl, MG, FT + 45000, int(SW * 0.75), 180000)
    ln(fl, "Water Corporation  |  EDAP Programme  |  March 2026", 8, LBLUE, first=True)

    fr = tx(sl, SW - MG - 550000, FT + 45000, 550000, 180000)
    ln(fr, f"{num} / 4", 8, LBLUE, align=PP_ALIGN.RIGHT, first=True)

def context_line(sl, text):
    c = tx(sl, MG, CT + 90000, SW - 2*MG, 260000)
    ln(c, text, 10.5, DARK, italic=True, first=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portfolio Vision
# ══════════════════════════════════════════════════════════════════════════════

def slide1(p):
    sl = blank(p)
    chrome(sl,
           "The EDAP Portfolio Vision",
           "The strategic intent driving Water Corporation\u2019s enterprise data transformation",
           "1")

    context_line(sl,
        "Water Corporation\u2019s data capability has grown organically across seven domains. "
        "EDAP exists to consolidate, govern, and unlock it.")

    # ── Vision quote box ──────────────────────────────────────────────────────
    VT, VH = CT + 390000, 1100000
    bx(sl, MG,        VT, SW - 2*MG, VH, NAVY)
    bx(sl, MG,        VT, 18000,     VH, TEAL)   # left accent stripe

    vq = tx(sl, MG + 60000, VT + PM, SW - 2*MG - 80000, VH - PM * 2)
    ln(vq,
       "\u201cTransition Water Corporation from fragmented, reactive data practices to a "
       "proactive, intelligence-driven utility \u2014 where compliance by design and "
       "data liquidity accelerate operational resilience and innovation.\u201d",
       14, WHITE, italic=True, first=True, sa=12)
    ln(vq,
       "EDAP Portfolio Vision  |  Water Corporation Data & Analytics Programme",
       9, LBLUE)

    # ── Three columns ─────────────────────────────────────────────────────────
    C3T  = VT + VH + 110000
    C3H  = FT - C3T - 80000
    GAP  = 60480
    CW   = (SW - 2*MG - 2*GAP) // 3

    cols = [
        ("The Challenge", [
            "\u2022  Six source systems in operational silos",
            "\u2022  Four pipeline technologies (Glue, ADF, SAP DS, SageMaker) uncoordinated",
            "\u2022  No unified governance \u2014 classification inconsistent, access ad hoc",
            "\u2022  Regulatory obligations met manually and reactively",
            "\u2022  Knowledge concentrated in individuals, not the platform",
        ]),
        ("The Aspiration", [
            "\u2022  Single governed platform across all seven business domains",
            "\u2022  New data source onboarded in days via configuration, not code",
            "\u2022  Compliance enforced architecturally at column level via ABAC",
            "\u2022  Certified data products with quality contracts and SLAs",
            "\u2022  Analytics and AI governed from experiment through to production",
        ]),
        ("How EDAP Delivers", [
            "\u2022  Databricks lakehouse + Unity Catalog across Dev / Test / Prod",
            "\u2022  Medallion architecture \u2014 10 structured zones from landing to gold",
            "\u2022  4-layer tagging: WAICP \u2192 sensitivity \u2192 access \u2192 operational metadata",
            "\u2022  9 epics, 95 features, SAFe PI delivery cadence",
            "\u2022  WC engineers certified, embedded, and self-sufficient at close",
        ]),
    ]

    for i, (title, bullets) in enumerate(cols):
        L = MG + i * (CW + GAP)
        bx(sl, L, C3T,  CW, C3H,    LGRAY, border=True)
        bx(sl, L, C3T,  CW, 310000, TEAL)

        ht = tx(sl, L + PM, C3T + 70000, CW - 2*PM, 200000)
        ln(ht, title, 11, WHITE, bold=True, first=True)

        bt = tx(sl, L + PM, C3T + 370000, CW - 2*PM, C3H - 420000)
        first_b = True
        for b in bullets:
            if first_b:
                ln(bt, b, 9.5, DARK, first=True, sa=3)
                first_b = False
            else:
                ln(bt, b, 9.5, DARK, sb=3, sa=3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Program Strategic Themes
# ══════════════════════════════════════════════════════════════════════════════

def slide2(p):
    sl = blank(p)
    chrome(sl,
           "Program Strategic Themes",
           "Five themes that define what the EDAP Programme delivers for Water Corporation",
           "2")

    context_line(sl,
        "Each theme maps directly to one or more delivery epics. "
        "Together they describe the transformation EDAP is contracted to deliver.")

    themes = [
        ("01", "Trusted Data Foundation",       "E1 \u00b7 E2 \u00b7 E3",
         "Build and operationalise a single governed lakehouse platform \u2014 replacing four fragmented pipeline "
         "technologies with a metadata-driven framework on Databricks. Medallion architecture across 10 structured "
         "zones, Unity Catalog for access governance, and certified data products with contractual quality guarantees."),

        ("02", "Compliance by Design",           "E4",
         "Embed regulatory compliance architecturally rather than managing it manually. WAICP, PRIS Act 2024, "
         "SOCI Act 2018, and State Records Act obligations enforced through tag-driven attribute-based access "
         "control (ABAC) and complete audit trails in Unity Catalog \u2014 demonstrable on demand."),

        ("03", "Discoverable and Liquid Data",   "E3 \u00b7 E6",
         "Make every data asset across all seven domains findable in Alation with full lineage from source "
         "through to Gold. Metadata stays current automatically via the OCF connector. Domain teams access "
         "and share data through governed, configuration-driven interfaces \u2014 not platform tickets."),

        ("04", "Intelligence at Scale",          "E7 \u00b7 E8",
         "Enable the full ML lifecycle on EDAP \u2014 experiment, train, register, deploy, and monitor models "
         "with governance from day one. Document and prove clear migration pathways from legacy platforms "
         "(Glue, ADF, SAP DS, SageMaker) so legacy workloads have a costed, proven path forward."),

        ("05", "Enduring Capability",            "E5 \u00b7 E9",
         "Deliver a self-sustaining platform and a self-sufficient WC team. DataOps excellence embeds CI/CD, "
         "automated testing, and non-prod cost guardrails into every pipeline. WC engineers are "
         "Databricks-certified, embedded hands-on through delivery, and own their domains at project close."),
    ]

    TH     = 840000
    GAP    = 52000
    START  = CT + 360000
    LW     = 1620000    # left label panel width

    for i, (num, name, epics, desc) in enumerate(themes):
        T = START + i * (TH + GAP)
        W = SW - 2*MG

        bx(sl, MG, T, W,  TH, LGRAY, border=True)
        bx(sl, MG, T, LW, TH, NAVY)

        # Number (large, teal)
        nt = tx(sl, MG + PM, T + 80000, 400000, 310000)
        ln(nt, num, 22, TEAL, bold=True, first=True)

        # Theme name
        nm = tx(sl, MG + PM, T + 400000, LW - 2*PM, 300000)
        ln(nm, name, 10.5, LBLUE, bold=True, first=True)

        # Epic refs
        ep = tx(sl, MG + PM, T + TH - 195000, LW - 2*PM, 170000)
        ln(ep, f"Epics: {epics}", 8.5, TEAL, first=True)

        # Description
        dt = tx(sl, MG + LW + PM, T + int(PM * 0.8), W - LW - 2*PM, TH - int(PM * 1.6))
        ln(dt, desc, 9.5, DARK, first=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Program Scope: Nine Epics
# ══════════════════════════════════════════════════════════════════════════════

def slide3(p):
    sl = blank(p)
    chrome(sl,
           "Program Scope \u2014 Nine Epics Across the Delivery",
           "95 features, sequenced foundation-first, spanning platform, governance, intelligence, and capability",
           "3")

    context_line(sl,
        "The nine epics are ordered architectural layers, not independent workstreams. "
        "E1 is the foundation for all that follows. E9 runs in parallel throughout delivery.")

    epics = [
        ("E1", "Platform Foundation",
         NAVY, "19 features",
         "EDAP environments operational, governed, and ready for development \u2014 "
         "the foundation everything else depends on"),

        ("E2", "Data Ingestion & Integration",
         TEAL, "14 features",
         "Data flowing reliably from source systems into EDAP Bronze layer via "
         "batch, CDC, streaming, and file-based patterns"),

        ("E3", "Transformation & Data Products",
         TEAL, "13 features",
         "Clean, conformed, business-ready data products in the Gold layer with "
         "dimensional models, semantic metrics, and quality contracts"),

        ("E4", "Security, Governance & Compliance",
         TEAL2, "16 features",
         "Data classified, protected, and compliant automatically \u2014 "
         "tag-driven ABAC, audit trails, and regulatory posture demonstrable on demand"),

        ("E5", "DataOps & Engineering Excellence",
         TEAL2, "12 features",
         "Automated delivery pipeline \u2014 code promotes safely, tests run automatically, "
         "non-prod costs stay at \u226410% of production"),

        ("E6", "Discovery & Catalogue",
         TEAL2, "7 features",
         "Every EDAP data asset discoverable in Alation with full metadata, "
         "lineage, and business context \u2014 kept current automatically"),

        ("E7", "Analytics & AI",
         NAVY, "6 features",
         "MLOps capabilities operational: experiment, train, deploy, and monitor models; "
         "three OMLIX algorithms running in the Customer domain"),

        ("E8", "Migration Pathways",
         NAVY, "6 features",
         "Every legacy platform (Glue, ADF, SAP DS, SageMaker) has a documented, "
         "costed, proven migration path to EDAP"),

        ("E9", "Capability Uplift & Advisory",
         RGBColor(0x1A, 0x3A, 0x2A), "17 features",
         "WC team self-sufficient on EDAP \u2014 trained, certified, embedded in delivery, "
         "with endorsed governance frameworks for data products and AI"),
    ]

    ROWS, COLS = 3, 3
    GAP  = 55000
    CW   = (SW - 2*MG - (COLS - 1)*GAP) // COLS
    ET   = CT + 370000
    RH   = (FT - ET - 60000 - (ROWS - 1)*GAP) // ROWS
    HBH  = 265000   # header band height inside each card

    for idx, (eid, ename, ecolor, nfeat, outcome) in enumerate(epics):
        row = idx // COLS
        col = idx % COLS
        L   = MG + col * (CW + GAP)
        T   = ET  + row * (RH + GAP)

        bx(sl, L, T, CW, RH,   LGRAY, border=True)
        bx(sl, L, T, CW, HBH,  ecolor)

        # Epic id + name
        ht = tx(sl, L + PM, T + 48000, CW - 2*PM - 300000, 200000)
        ln(ht, f"{eid}  {ename}", 10.5, WHITE, bold=True, first=True)

        # Feature count (top-right)
        fc = tx(sl, L + CW - PM - 380000, T + 70000, 380000, 160000)
        ln(fc, nfeat, 8.5, LBLUE, align=PP_ALIGN.RIGHT, first=True)

        # Outcome text
        ot = tx(sl, L + PM, T + HBH + 60000, CW - 2*PM, RH - HBH - 100000)
        ln(ot, outcome, 9.5, DARK, first=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Project Outcomes
# ══════════════════════════════════════════════════════════════════════════════

def slide4(p):
    sl = blank(p)
    chrome(sl,
           "Project Outcomes \u2014 Supporting the Strategic Vision",
           "How this project\u2019s deliverables connect to each of the five strategic themes",
           "4")

    context_line(sl,
        "Each outcome is traceable to a scope item and an epic. "
        "Together they confirm EDAP has landed \u2014 not just been deployed.")

    outcomes = [
        ("Trusted Data\nFoundation",     "E1 \u00b7 E2 \u00b7 E3", [
            "Databricks lakehouse across Dev / Test / Prod with Unity Catalog metastore and domain-isolated catalogues",
            "Metadata-driven ingestion framework: configuration-based onboarding for batch, CDC, streaming, and file sources",
            "Business-ready dimensional models in Gold layer with contractual quality guarantees and automated SLA tracking",
        ]),
        ("Compliance\nby Design",        "E4", [
            "4-layer tagging model live: WAICP classification \u2192 sensitivity reason \u2192 access handling \u2192 operational metadata",
            "Alation\u2013Unity Catalog bidirectional sync: steward classifications automatically enforce ABAC policies in Databricks",
            "PRIS / SOCI / State Records compliance demonstrable on demand from Unity Catalog system tables",
        ]),
        ("Discoverable\nand Liquid Data","E3 \u00b7 E6", [
            "Every EDAP asset discoverable in Alation with full end-to-end lineage from source system through to Gold",
            "Metadata current automatically via OCF connector; stewardship workflows operational with clear accountability",
        ]),
        ("Intelligence\nat Scale",       "E7 \u00b7 E8", [
            "MLOps pipeline operational: experiment tracking, feature store, model registry, deployment, and drift monitoring",
            "Three OMLIX algorithms in production; migration pathways for Glue, ADF, SAP DS, and SageMaker documented and proven",
        ]),
        ("Enduring\nCapability",         "E5 \u00b7 E9", [
            "DataOps CI/CD: branching, automated test pyramid, one-click promotion Dev \u2192 Test \u2192 Prod, non-prod cost \u226410% of production",
            "WC engineers Databricks-certified, embedded in code review, and leading delivery of their own domain pipelines",
            "Reusable runbooks, templates, and frameworks so WC owns and operates the platform beyond project close",
        ]),
    ]

    LW    = 1640000   # left label panel width
    GAP   = 45000
    START = CT + 370000
    AVAIL = FT - START - 60000

    # Calculate heights proportional to bullet count
    raw_h = [260000 + len(b) * 230000 for _, _, b in outcomes]
    total = sum(raw_h) + (len(raw_h) - 1) * GAP
    scale = min(1.0, AVAIL / total)
    heights = [int(h * scale) for h in raw_h]

    T = START
    for i, ((theme, epics_ref, bullets), H) in enumerate(zip(outcomes, heights)):
        W = SW - 2*MG

        bx(sl, MG, T, W,  H, LGRAY, border=True)
        bx(sl, MG, T, LW, H, NAVY)

        # Theme name (multi-line in label)
        tm = tx(sl, MG + PM, T + int(PM * 0.6), LW - 2*PM, H - int(PM * 1.2))
        lines = theme.split("\n")
        ln(tm, lines[0], 10.5, WHITE, bold=True, first=True)
        if len(lines) > 1:
            ln(tm, lines[1], 10.5, WHITE, bold=True)
        ln(tm, f"Epics: {epics_ref}", 8.5, TEAL, sb=5)

        # Bullet outcomes
        bt = tx(sl, MG + LW + PM, T + int(PM * 0.6), W - LW - 2*PM, H - int(PM * 1.2))
        first_b = True
        for b in bullets:
            if first_b:
                ln(bt, f"\u2022  {b}", 9, DARK, first=True, sa=3)
                first_b = False
            else:
                ln(bt, f"\u2022  {b}", 9, DARK, sb=4, sa=3)

        T += H + GAP


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    prs = mk()
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    out = "/home/user/DataWikis/delivery/edap-executive-slide-pack.pptx"
    prs.save(out)
    print(f"Saved → {out}")
